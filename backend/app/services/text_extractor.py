"""Extract plain text from various document formats.

Supports:
  - Text-based: .txt, .md, .csv, .html, .htm  (decode bytes directly)
  - Word:       .docx                           (python-docx)
  - PDF:        .pdf                            (PyMuPDF / fitz)
  - Excel:      .xlsx                           (openpyxl)
  - PowerPoint: .pptx                           (python-pptx)

For unsupported or unreadable formats the extractor returns an empty string
so the caller can decide how to proceed (e.g. skip glossary generation).
"""

import io
import logging

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Text-based formats – just decode the bytes
# ---------------------------------------------------------------------------

_TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".html", ".htm", ".tsv", ".json", ".xml"}


def _extract_text_plain(data: bytes) -> str:
    """Decode raw bytes as UTF-8 (with fallback to latin-1)."""
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1", errors="replace")


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def _extract_text_docx(data: bytes) -> str:
    try:
        from docx import Document  # python-docx

        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except Exception as exc:
        logger.warning("DOCX text extraction failed: %s", exc)
        return ""


# ---------------------------------------------------------------------------
# PDF  (PyMuPDF)
# ---------------------------------------------------------------------------

def _extract_text_pdf(data: bytes) -> str:
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=data, filetype="pdf")
        pages = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return "\n".join(pages)
    except Exception as exc:
        logger.warning("PDF text extraction failed: %s", exc)
        return ""


# ---------------------------------------------------------------------------
# XLSX  (openpyxl)
# ---------------------------------------------------------------------------

def _extract_text_xlsx(data: bytes) -> str:
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines: list[str] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append("\t".join(cells))
        wb.close()
        return "\n".join(lines)
    except Exception as exc:
        logger.warning("XLSX text extraction failed: %s", exc)
        return ""


# ---------------------------------------------------------------------------
# PPTX  (python-pptx)
# ---------------------------------------------------------------------------

def _extract_text_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
        return "\n".join(texts)
    except Exception as exc:
        logger.warning("PPTX text extraction failed: %s", exc)
        return ""


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

_EXTRACTORS = {
    ".docx": _extract_text_docx,
    ".pdf": _extract_text_pdf,
    ".xlsx": _extract_text_xlsx,
    ".pptx": _extract_text_pptx,
}


def extract_text(data: bytes, file_extension: str) -> str:
    """Extract plain text from *data* based on *file_extension*.

    Returns an empty string if the format is unsupported or extraction fails.
    """
    ext = file_extension.lower()

    if ext in _TEXT_EXTENSIONS:
        return _extract_text_plain(data)

    extractor = _EXTRACTORS.get(ext)
    if extractor:
        return extractor(data)

    logger.info("No text extractor for extension '%s'", ext)
    return ""
